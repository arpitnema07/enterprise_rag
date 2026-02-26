"""
PPTX Extraction Module

Extracts text, tables, speaker notes, and embedded images from PowerPoint (.pptx) files.
Legacy .ppt files should be converted to .pdf first.
Each slide is processed to form semantic context blocks.
"""

import os
import logging
from typing import List, Dict, Any
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import base64
import subprocess

logger = logging.getLogger(__name__)


def extract_pptx_content(pptx_path: str) -> List[Dict[str, Any]]:
    """
    Extracts content from a .pptx file, slide by slide.

    Args:
        pptx_path: Path to the .pptx file.

    Returns:
        List of dicts, where each dict represents a slide with:
        - 'page_number': Slide index (1-based)
        - 'text': Combined extracted text, table data, and notes
        - 'has_tables': Boolean indicating if tables were found
        - 'extraction_method': String indicating pptx parser
        - 'images': List of extracted image bytes (base64 encoded)
    """
    logger.info(f"Extracting PPTX content from {pptx_path}")

    try:
        prs = Presentation(pptx_path)
    except Exception as e:
        logger.error(f"Failed to load PPTX {pptx_path}: {e}")
        return []

    slides_data = []

    # Optional: If you want to use vision OCR on images, import here
    # from .vision import process_image_with_vision # This line was moved inside the loop as per the change

    for slide_idx, slide in enumerate(prs.slides):
        slide_num = slide_idx + 1
        slide_texts = []
        has_tables = False
        images_b64 = []

        # 1. Extract Speaker Notes
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                slide_texts.append(f"--- Speaker Notes ---\n{notes}\n--- End Notes ---")

        # 2. Extract Shapes (TextFrames, Tables, Pictures)
        for shape in slide.shapes:
            # Text Boxes
            if shape.has_text_frame:
                text = shape.text.strip()
                if text:
                    slide_texts.append(text)

            # Tables
            elif shape.has_table:
                has_tables = True
                table_md = _extract_table_as_markdown(shape.table)
                if table_md:
                    slide_texts.append(
                        f"--- Table Data ---\n{table_md}\n--- End Table ---"
                    )

            # Pictures
            elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    img_bytes = shape.image.blob
                    b64_img = base64.b64encode(img_bytes).decode("utf-8")
                    images_b64.append(
                        {"image_b64": b64_img, "image_index": len(images_b64)}
                    )

                    # Optional: Process image with vision LLM immediately to get caption
                    try:
                        from .vision import caption_image

                        caption = caption_image(b64_img)
                        if caption:
                            slide_texts.append(
                                f"--- Image Content ---\n{caption}\n--- End Image ---"
                            )
                    except Exception as ve:
                        logger.warning(
                            f"Vision OCR failed for image on slide {slide_num}: {ve}"
                        )

                except Exception as e:
                    logger.warning(f"Failed to extract image on slide {slide_num}: {e}")

        # Combine all text for the slide
        combined_text = "\n\n".join(slide_texts)

        slides_data.append(
            {
                "page_number": slide_num,
                "text": combined_text,
                "has_tables": has_tables,
                "extraction_method": "python-pptx",
                "images": images_b64,
            }
        )

    return slides_data


def _extract_table_as_markdown(table) -> str:
    """Helper to convert python-pptx table object to Markdown."""
    if not table.columns or not table.rows:
        return ""

    rows_data = []
    num_cols = len(table.columns)

    for row in table.rows:
        row_cells = []
        for cell in row.cells:
            # Clean newlines inside cells
            cell_text = cell.text.replace("\n", " ").replace("\r", "").strip()
            row_cells.append(cell_text)
        rows_data.append(row_cells)

    if not rows_data:
        return ""

    md_lines = []

    # Header
    header = rows_data[0]
    md_lines.append("| " + " | ".join(header) + " |")

    # Separator
    separator = ["---"] * num_cols
    md_lines.append("| " + " | ".join(separator) + " |")

    # Body rows
    for row in rows_data[1:]:
        md_lines.append("| " + " | ".join(row) + " |")

    return "\n".join(md_lines)


def convert_ppt_to_pdf(ppt_path: str, output_dir: str = None) -> str:
    """
    Converts a legacy .ppt file to .pdf using LibreOffice headless.

    Args:
        ppt_path: Path to the .ppt file
        output_dir: Optional directory to save the PDF. If None, uses the same directory.

    Returns:
        Path to the generated .pdf file, or None if conversion failed.
    """
    if not output_dir:
        output_dir = os.path.dirname(ppt_path)

    logger.info(f"Converting legacy PPT to PDF: {ppt_path} -> {output_dir}")

    try:
        # Command to convert using LibreOffice
        # Requires libreoffice to be installed and available in PATH
        cmd = [
            "soffice",
            "--headless",
            "--convert-to",
            "pdf",
            "--outdir",
            output_dir,
            ppt_path,
        ]

        subprocess.run(
            cmd, stdout=subprocess.PIPE, stderr=subprocess.PIPE, text=True, check=True
        )

        # Determine the expected output filename
        base_name = os.path.splitext(os.path.basename(ppt_path))[0]
        expected_pdf_path = os.path.join(output_dir, f"{base_name}.pdf")

        if os.path.exists(expected_pdf_path):
            logger.info(f"Successfully converted PPT to PDF: {expected_pdf_path}")
            return expected_pdf_path
        else:
            logger.error(
                f"LibreOffice command succeeded but expected output file not found: {expected_pdf_path}"
            )
            return None

    except subprocess.CalledProcessError as e:
        logger.error(f"LibreOffice conversion failed: {e.stderr}")
        return None
    except Exception as e:
        logger.error(f"Error invoking LibreOffice for PPT conversion: {e}")
        return None
